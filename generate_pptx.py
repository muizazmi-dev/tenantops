from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0b, 0x2d, 0x57)
AZURE      = RGBColor(0x00, 0x78, 0xd4)
LIGHT_BLUE = RGBColor(0xe8, 0xf4, 0xfd)
WHITE      = RGBColor(0xff, 0xff, 0xff)
ORANGE     = RGBColor(0xf0, 0x8a, 0x24)
DARK_TEXT  = RGBColor(0x1a, 0x1a, 0x2e)
MID_GREY   = RGBColor(0xcc, 0xd6, 0xe0)
CODE_BG    = RGBColor(0x1e, 0x1e, 0x2e)
CODE_FG    = RGBColor(0xa8, 0xd8, 0xea)
GREEN      = RGBColor(0x10, 0x7c, 0x10)
PURPLE     = RGBColor(0x6b, 0x37, 0x9c)
RED        = RGBColor(0xc4, 0x3f, 0x12)
LIGHT_GREY = RGBColor(0xf5, 0xf7, 0xfa)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ── Helpers ────────────────────────────────────────────────────────────────
def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h, font_size=Pt(12), bold=False,
                 italic=False, color=DARK_TEXT, align=PP_ALIGN.LEFT,
                 wrap=True, font_name="Segoe UI"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb

def header_bar(slide, title, subtitle=None, bar_color=NAVY,
               title_color=WHITE, sub_color=ORANGE):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.18), fill=bar_color)
    add_text_box(slide, title, Inches(0.35), Inches(0.1), Inches(9), Inches(0.6),
                 font_size=Pt(26), bold=True, color=title_color)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.35), Inches(0.68), Inches(10), Inches(0.42),
                     font_size=Pt(14), color=sub_color, italic=True)

def layer_badge(slide, label, x=Inches(10.9), y=Inches(0.12)):
    add_rect(slide, x, y, Inches(2.1), Inches(0.42), fill=ORANGE)
    add_text_box(slide, label, x + Inches(0.08), y + Inches(0.05),
                 Inches(1.95), Inches(0.35),
                 font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def content_box(slide, x, y, w, h, fill=LIGHT_BLUE, line=MID_GREY):
    return add_rect(slide, x, y, w, h, fill=fill, line=line, line_w=Pt(1))

def code_block(slide, code_lines, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=CODE_BG, line=AZURE, line_w=Pt(1.2))
    txb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.1),
                                    w - Inches(0.24), h - Inches(0.2))
    tf = txb.text_frame; tf.word_wrap = False
    first = True
    for line in code_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9.5)
        run.font.name = "Cascadia Code"
        run.font.color.rgb = CODE_FG

def harrow(slide, x, y, w, color=NAVY):
    add_rect(slide, x, y - Inches(0.022), w - Inches(0.07), Inches(0.044), fill=color)
    add_rect(slide, x + w - Inches(0.09), y - Inches(0.065), Inches(0.09), Inches(0.13), fill=color)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, 0, Inches(4.8), SLIDE_W, Inches(0.06), fill=ORANGE)
add_text_box(s, "TenantOps", Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
             font_size=Pt(56), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "Multi-Tenant SaaS Demo — Architecture Walkthrough",
             Inches(1), Inches(2.5), Inches(11.3), Inches(0.7),
             font_size=Pt(22), color=ORANGE, align=PP_ALIGN.CENTER)
add_text_box(s,
    "Client  ·  Edge / WAF  ·  API Gateway  ·  Application  ·  Service Layer\n"
    "Data  ·  Infrastructure  ·  AI",
    Inches(1), Inches(3.3), Inches(11.3), Inches(0.9),
    font_size=Pt(14), color=MID_GREY, align=PP_ALIGN.CENTER)
add_rect(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.04), fill=AZURE)
add_text_box(s,
    "Defense-in-depth multi-tenant isolation  ·  Semantic Kernel RAG  ·  Azure-native IaC",
    Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
    font_size=Pt(11), color=MID_GREY, align=PP_ALIGN.CENTER, italic=True)

add_notes(s, """SLIDE 1 — TITLE (Opening hook ~60 seconds)

SAY: "TenantOps is a production-grade multi-tenant SaaS demo built on Azure. It shows how to build a system where tenant data isolation is enforced at every layer — not just the database — but also at the browser, edge, API gateway, and AI pipeline.

We have two demo tenants: Contoso HVAC and Fabrikam Services. We will prove that a Contoso user cannot see Fabrikam data, even if they forge HTTP headers, obtain a network trace, or directly call the APIs.

Key stats:
- 5 independent isolation enforcement points (WAF, APIM, Middleware, SQL RLS, AI Search filter)
- 4 ASP.NET Core 8 microservices running in Azure Container Apps
- Next.js 14 frontend with React Server Components and httpOnly JWT cookies
- Semantic Kernel RAG pipeline with tenant-isolated vector search
- Full Terraform IaC: one command deploys to Azure

BEFORE STARTING: Confirm docker compose is running and make migrate + make seed have completed.
Check: docker compose ps  — all 14 containers should show 'healthy'
Confirm: localhost:3000, localhost:5001, 5002, 5003, 5004 are all responding.""")

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Full Architecture Diagram
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xf0, 0xf4, 0xf8))
header_bar(s, "Full System Architecture Diagram",
           "Browser → Front Door WAF → APIM → Container Apps → Azure SQL + AI Search + OpenAI")

# Zone definitions: (label, x, width, fill_color)
ZONE_Y = 1.28; ZONE_H = 6.08
zones = [
    ("INTERNET",         0.15,  1.35, RGBColor(0xe5, 0xe8, 0xf5)),
    ("EDGE / CDN / WAF", 1.68,  1.80, RGBColor(0xd5, 0xe5, 0xf8)),
    ("API GATEWAY",      3.66,  1.68, RGBColor(0xe5, 0xd5, 0xf8)),
    ("APPLICATION",      5.52,  3.48, RGBColor(0xd5, 0xf0, 0xe0)),
    ("DATA & AI",        9.18,  3.98, RGBColor(0xf8, 0xee, 0xd5)),
]
for zlabel, zx, zw, zfill in zones:
    add_rect(s, Inches(zx), Inches(ZONE_Y), Inches(zw), Inches(ZONE_H),
             fill=zfill, line=MID_GREY, line_w=Pt(0.8))
    add_text_box(s, zlabel, Inches(zx + 0.06), Inches(ZONE_Y + 0.06),
                 Inches(zw - 0.12), Inches(0.28),
                 font_size=Pt(8), bold=True, color=RGBColor(0x55, 0x55, 0x88),
                 align=PP_ALIGN.CENTER)

# ── Internet: Browser ──
BRW_X = 0.22; BRW_Y = 3.55
add_rect(s, Inches(BRW_X), Inches(BRW_Y), Inches(1.18), Inches(1.0),
         fill=NAVY, line=WHITE, line_w=Pt(1.2))
add_text_box(s, "Browser\n(User)", Inches(BRW_X), Inches(BRW_Y + 0.1),
             Inches(1.18), Inches(0.85),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "localhost:3000", Inches(BRW_X), Inches(BRW_Y + 1.08),
             Inches(1.2), Inches(0.22),
             font_size=Pt(7.5), color=AZURE, align=PP_ALIGN.CENTER,
             font_name="Cascadia Code")

# ── Edge: Front Door box ──
FD_X = 1.76; FD_Y = 1.82; FD_W = 1.64; FD_H = 5.22
add_rect(s, Inches(FD_X), Inches(FD_Y), Inches(FD_W), Inches(FD_H),
         fill=RGBColor(0x00, 0x62, 0xb1), line=WHITE, line_w=Pt(1.2))
add_text_box(s, "Azure Front Door\nPremium", Inches(FD_X), Inches(FD_Y + 0.08),
             Inches(FD_W), Inches(0.55),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

fd_pills = [
    ("WAF  Prevention Mode\nOWASP 2.1 + Bot Manager 1.1", RED),
    ("600 req/min per IP\nDDoS edge mitigation", RGBColor(0x8b, 0x00, 0x00)),
    ("Route /*  → Web App", RGBColor(0x00, 0x48, 0x8e)),
    ("Route /api/* → APIM", RGBColor(0x50, 0x20, 0x80)),
    ("CDN + TLS Termination", RGBColor(0x00, 0x4c, 0x8c)),
]
pilly = FD_Y + 0.72
for ptxt, pcol in fd_pills:
    add_rect(s, Inches(FD_X + 0.1), Inches(pilly), Inches(FD_W - 0.2), Inches(0.5),
             fill=pcol, line=WHITE, line_w=Pt(0.3))
    add_text_box(s, ptxt, Inches(FD_X + 0.1), Inches(pilly + 0.04),
                 Inches(FD_W - 0.2), Inches(0.44),
                 font_size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    pilly += 0.62

# ── Gateway: APIM box ──
AP_X = 3.74; AP_Y = 1.82; AP_W = 1.58; AP_H = 5.22
add_rect(s, Inches(AP_X), Inches(AP_Y), Inches(AP_W), Inches(AP_H),
         fill=PURPLE, line=WHITE, line_w=Pt(1.2))
add_text_box(s, "Azure APIM", Inches(AP_X), Inches(AP_Y + 0.08),
             Inches(AP_W), Inches(0.42),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

apim_pills = [
    ("Strip x-tenant-id", RGBColor(0x99, 0x00, 0x00)),
    ("Validate JWT\n(Entra or local key)", RGBColor(0x00, 0x44, 0x88)),
    ("Inject x-tenant-id\nfrom JWT claim", RGBColor(0x00, 0x66, 0x44)),
    ("Rate limit 60/min\nper tenant", RGBColor(0x55, 0x22, 0x88)),
    ("/tenant  /core\n/ai  /identity", RGBColor(0x33, 0x33, 0x77)),
]
py_ap = AP_Y + 0.6
for atxt, acol in apim_pills:
    add_rect(s, Inches(AP_X + 0.1), Inches(py_ap), Inches(AP_W - 0.2), Inches(0.5),
             fill=acol, line=WHITE, line_w=Pt(0.3))
    add_text_box(s, atxt, Inches(AP_X + 0.1), Inches(py_ap + 0.04),
                 Inches(AP_W - 0.2), Inches(0.44),
                 font_size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    py_ap += 0.62

# ── Application: Next.js + Container Apps ──
APP_X = 5.60; APP_Y = 1.82; APP_W = 3.32

# Next.js bar
add_rect(s, Inches(APP_X), Inches(APP_Y), Inches(APP_W), Inches(0.95),
         fill=AZURE, line=WHITE, line_w=Pt(1.2))
add_text_box(s, "Next.js 14  (Port 3000)",
             Inches(APP_X), Inches(APP_Y + 0.05), Inches(APP_W), Inches(0.4),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "RSC  ·  Edge Middleware  ·  BFF /api/chat",
             Inches(APP_X), Inches(APP_Y + 0.48), Inches(APP_W), Inches(0.35),
             font_size=Pt(8), color=RGBColor(0xcc, 0xee, 0xff), align=PP_ALIGN.CENTER)

# Container Apps environment
CA_Y = APP_Y + 1.08; CA_H = 4.0
add_rect(s, Inches(APP_X), Inches(CA_Y), Inches(APP_W), Inches(CA_H),
         fill=RGBColor(0x09, 0x24, 0x44), line=ORANGE, line_w=Pt(1.5))
add_text_box(s, "Azure Container Apps Environment",
             Inches(APP_X + 0.05), Inches(CA_Y + 0.07), Inches(APP_W - 0.1), Inches(0.28),
             font_size=Pt(8), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

svcs = [
    ("identity-api    :5004   JWT issuance + PBKDF2",   AZURE),
    ("tenant-api      :5001   slug → tenantId resolve",  RGBColor(0x00, 0x4c, 0x8c)),
    ("core-api        :5002   tickets + documents",      GREEN),
    ("ai-orchestrator :5003   RAG embed+search+chat",    RED),
]
sy = CA_Y + 0.46
for stxt, scol in svcs:
    add_rect(s, Inches(APP_X + 0.1), Inches(sy), Inches(APP_W - 0.2), Inches(0.65),
             fill=scol, line=WHITE, line_w=Pt(0.5))
    add_text_box(s, stxt, Inches(APP_X + 0.15), Inches(sy + 0.1),
                 Inches(APP_W - 0.3), Inches(0.5),
                 font_size=Pt(8.5), bold=True, color=WHITE, font_name="Cascadia Code")
    sy += 0.80

# Middleware footer inside CA
add_rect(s, Inches(APP_X + 0.1), Inches(sy + 0.05), Inches(APP_W - 0.2), Inches(0.42),
         fill=ORANGE, line=WHITE, line_w=Pt(0.5))
add_text_box(s, "TenantResolutionMiddleware  (all 4 services)",
             Inches(APP_X + 0.12), Inches(sy + 0.1), Inches(APP_W - 0.26), Inches(0.34),
             font_size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Data & AI ──
D_X = 9.26; D_Y = 1.82; D_W = 3.82

# SQL
add_rect(s, Inches(D_X), Inches(D_Y), Inches(D_W), Inches(1.6),
         fill=GREEN, line=WHITE, line_w=Pt(1.2))
add_text_box(s, "Azure SQL Database",
             Inches(D_X), Inches(D_Y + 0.08), Inches(D_W), Inches(0.38),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "Row-Level Security  ·  fn_tenant_predicate\nSESSION_CONTEXT(TenantId)  ·  FILTER + BLOCK\napp.Tenants / Users / Tickets / Documents",
             Inches(D_X + 0.1), Inches(D_Y + 0.5), Inches(D_W - 0.2), Inches(1.0),
             font_size=Pt(8), color=RGBColor(0xd0, 0xff, 0xd0), align=PP_ALIGN.CENTER)

# AI Search
add_rect(s, Inches(D_X), Inches(D_Y + 1.75), Inches(D_W), Inches(1.6),
         fill=RGBColor(0x00, 0x62, 0xb1), line=WHITE, line_w=Pt(1.2))
add_text_box(s, "Azure AI Search",
             Inches(D_X), Inches(D_Y + 1.83), Inches(D_W), Inches(0.38),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "Hybrid BM25 + KNN  ·  1536-dim vectors\ntenantId field + OData filter mandatory\nTop-5 chunks → grounded prompt context",
             Inches(D_X + 0.1), Inches(D_Y + 2.25), Inches(D_W - 0.2), Inches(1.0),
             font_size=Pt(8), color=RGBColor(0xcc, 0xee, 0xff), align=PP_ALIGN.CENTER)

# OpenAI
add_rect(s, Inches(D_X), Inches(D_Y + 3.50), Inches(D_W), Inches(1.88),
         fill=RED, line=WHITE, line_w=Pt(1.2))
add_text_box(s, "Azure OpenAI",
             Inches(D_X), Inches(D_Y + 3.58), Inches(D_W), Inches(0.38),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "gpt-4o-mini  — Chat Completion\ntext-embedding-3-small  — 1536-dim\nGrounded to context only  ·  Structured citations",
             Inches(D_X + 0.1), Inches(D_Y + 4.0), Inches(D_W - 0.2), Inches(1.2),
             font_size=Pt(8), color=RGBColor(0xff, 0xcc, 0xaa), align=PP_ALIGN.CENTER)

# ── Arrows ──
AC = RGBColor(0x33, 0x44, 0x88)
# Browser → Front Door (center of browser ~y=4.05)
harrow(s, Inches(BRW_X + 1.18), Inches(4.05), Inches(0.40), AC)
add_text_box(s, "HTTPS", Inches(BRW_X + 1.18), Inches(3.86), Inches(0.4), Inches(0.2),
             font_size=Pt(7), color=AC, align=PP_ALIGN.CENTER)

# FD → Next.js (Next.js center ~y=2.30)
harrow(s, Inches(FD_X + FD_W), Inches(2.30), Inches(0.30), AC)
add_text_box(s, "Web/*", Inches(FD_X + FD_W + 0.01), Inches(2.13), Inches(0.3), Inches(0.2),
             font_size=Pt(7), color=AC, align=PP_ALIGN.CENTER)

# FD → APIM (/api/* route, center ~y=3.5)
harrow(s, Inches(FD_X + FD_W), Inches(3.50), Inches(0.22), AC)
add_text_box(s, "/api/*", Inches(FD_X + FD_W + 0.01), Inches(3.32), Inches(0.3), Inches(0.2),
             font_size=Pt(7), color=AC, align=PP_ALIGN.CENTER)

# APIM → Services (center of CA ~y=4.0)
harrow(s, Inches(AP_X + AP_W), Inches(4.0), Inches(0.25), AC)
add_text_box(s, "JWT\nverif.", Inches(AP_X + AP_W + 0.01), Inches(3.72), Inches(0.26), Inches(0.32),
             font_size=Pt(7), color=AC, align=PP_ALIGN.CENTER)

# App → SQL (all services; SQL center ~y=2.62)
harrow(s, Inches(APP_X + APP_W), Inches(2.62), Inches(0.30), AC)
add_text_box(s, "SQL+RLS", Inches(APP_X + APP_W + 0.01), Inches(2.44), Inches(0.3), Inches(0.2),
             font_size=Pt(7), color=AC, align=PP_ALIGN.CENTER)

# App → AI Search (ai-orchestrator → Search; Search center ~y=3.57)
harrow(s, Inches(APP_X + APP_W), Inches(3.57), Inches(0.30), AC)
add_text_box(s, "search", Inches(APP_X + APP_W + 0.01), Inches(3.39), Inches(0.3), Inches(0.2),
             font_size=Pt(7), color=AC, align=PP_ALIGN.CENTER)

# App → OpenAI (ai-orchestrator → OpenAI; OpenAI center ~y=5.76)
harrow(s, Inches(APP_X + APP_W), Inches(5.76), Inches(0.30), AC)
add_text_box(s, "embed\n+chat", Inches(APP_X + APP_W + 0.01), Inches(5.48), Inches(0.3), Inches(0.3),
             font_size=Pt(7), color=AC, align=PP_ALIGN.CENTER)

add_notes(s, """SLIDE 2 — FULL ARCHITECTURE DIAGRAM (2-3 minutes)

Walk through left to right:

INTERNET ZONE: "The browser is the only public-facing client. It connects to one URL — Azure Front Door. No service is directly reachable from the internet."

EDGE / CDN / WAF: "Azure Front Door Premium is the single ingress point. Before any request reaches a service:
- WAF applies OWASP 2.1 managed rule set in Prevention mode — malicious requests are blocked, not just logged.
- Bot Manager 1.1 blocks known bad bots.
- Rate limiting: 600 req/min per IP — DDoS protection at the edge, before APIM.
- TLS is terminated here. All HTTP is redirected to HTTPS.
- Two routes: /* goes to the Next.js web app; /api/* goes to APIM."

API GATEWAY (APIM): "Every API call is intercepted by APIM. The policies run in order:
1. Strip x-tenant-id — clients CANNOT set their own tenant header.
2. Validate JWT — checks signature, expiry, issuer (Entra ID in prod, local symmetric key in dev).
3. Inject x-tenant-id from the JWT's tid_app claim — the authoritative tenant is now in the header.
4. Rate limit 60 req/min per tenant — per-tenant throttling.
Even if a developer forgets [Authorize] on a .NET endpoint, APIM blocks unauthenticated calls."

APPLICATION ZONE: "Four microservices in Azure Container Apps. Each one shares TenantResolutionMiddleware from the shared library, which re-validates tenant claims. Next.js runs separately and uses a BFF pattern for the AI chat endpoint."

DATA & AI: "Three backend stores:
- Azure SQL with Row-Level Security — every SELECT/INSERT automatically filters by SESSION_CONTEXT(TenantId).
- Azure AI Search — every indexed chunk has a tenantId field; every search query applies an OData filter.
- Azure OpenAI — gpt-4o-mini for chat, text-embedding-3-small for vector indexing.

The arrows show logical data flow. All services write to SQL. Only ai-orchestrator talks to AI Search and OpenAI." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Overview (layer summary)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Architecture Overview",
           "8 independent layers — every layer independently enforces tenant boundaries")

LAYERS = [
    ("Client\n(Next.js 14)", RGBColor(0x00, 0x78, 0xd4)),
    ("Edge\n(Front Door\n+ WAF)",  RGBColor(0x00, 0x62, 0xb1)),
    ("API Gateway\n(APIM)",        RGBColor(0x6b, 0x37, 0x9c)),
    ("Application\n(4 × .NET 8)",  RGBColor(0x0b, 0x2d, 0x57)),
    ("Service Layer\n(Shared Lib)",RGBColor(0x10, 0x3a, 0x6b)),
    ("Data\n(SQL + RLS)",          RGBColor(0x10, 0x7c, 0x10)),
    ("Infra\n(Terraform\n+Docker)",RGBColor(0x6b, 0x37, 0x9c)),
    ("AI\n(SK + OpenAI\n+Search)", RGBColor(0xc4, 0x3f, 0x12)),
]
bw = Inches(1.42); bh = Inches(1.6); by = Inches(2.5); gap = Inches(0.12); start_x = Inches(0.35)
for i, (label, col) in enumerate(LAYERS):
    bx = start_x + i * (bw + gap)
    add_rect(s, bx, by, bw, bh, fill=col)
    add_text_box(s, label, bx + Inches(0.06), by + Inches(0.25),
                 bw - Inches(0.12), bh - Inches(0.3),
                 font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(LAYERS) - 1:
        add_rect(s, bx + bw + Inches(0.01), by + bh / 2 - Inches(0.03),
                 gap - Inches(0.01), Inches(0.06), fill=NAVY)

stats = [
    ("2 Demo Tenants", "Contoso HVAC · Fabrikam Services"),
    ("4 Microservices", "identity · tenant · core · ai-orchestrator"),
    ("5 Isolation Points", "WAF · APIM · Middleware · RLS · AI Filter"),
    ("RAG Pipeline", "Embed → Search → Ground → Cite"),
]
sy = Inches(4.5); sw = Inches(3.1)
for i, (h, d) in enumerate(stats):
    sx = Inches(0.35) + i * (sw + Inches(0.18))
    content_box(s, sx, sy, sw, Inches(0.9))
    add_text_box(s, h, sx + Inches(0.1), sy + Inches(0.05), sw - Inches(0.2), Inches(0.38),
                 font_size=Pt(11), bold=True, color=NAVY)
    add_text_box(s, d, sx + Inches(0.1), sy + Inches(0.42), sw - Inches(0.2), Inches(0.4),
                 font_size=Pt(9), color=DARK_TEXT, italic=True)

add_text_box(s, "Prerequisites:  docker compose up -d  →  make migrate  →  make seed",
             Inches(0.35), Inches(6.8), Inches(12.6), Inches(0.4),
             font_size=Pt(10), color=AZURE, italic=True, font_name="Cascadia Code")

add_notes(s, """SLIDE 3 — ARCHITECTURE OVERVIEW (1-2 minutes)

SAY: "Before we dive into each layer, here's the big picture. There are 8 layers and 5 independent isolation enforcement points. The word 'independent' is important: even if one layer fails or is misconfigured, the other four still prevent cross-tenant access.

Walk through each box left to right:
1. Client — Next.js 14 with React Server Components. JWT stored in httpOnly cookie, never in localStorage.
2. Edge — Azure Front Door + WAF. First line of defense. Never directly reachable backend.
3. API Gateway — APIM validates every JWT and strips/re-injects tenant headers.
4. Application — 4 ASP.NET Core 8 microservices, each with [Authorize] and middleware.
5. Service Layer — Shared middleware that validates tenant context on every request.
6. Data — SQL Row-Level Security automatically filters all queries by TenantId.
7. Infra — Terraform IaC, Docker locally. Managed identities in Azure (no secrets in code).
8. AI — Semantic Kernel RAG with tenant-isolated vector search in Azure AI Search.

BEFORE DEMO: Run 'docker compose ps' and show all 14 containers healthy. If any are unhealthy, run 'make restart'." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Layer 1: Client
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 1 — Client Layer",
           "Next.js 14  ·  React Server Components  ·  Edge Middleware  ·  httpOnly JWT cookies")
layer_badge(s, "CLIENT LAYER")

content_box(s, Inches(0.3), Inches(1.3), Inches(4.1), Inches(5.9))
add_text_box(s, "Pages & Routes", Inches(0.45), Inches(1.38), Inches(3.8), Inches(0.4),
             font_size=Pt(13), bold=True, color=NAVY)
pages = [
    "/  — Home: tenant picker cards",
    "/login  — Server action → identity-api → httpOnly cookie",
    "/admin  — Platform admin tenant list",
    "/t/[slug]  — Tenant home  (RSC, branding from ThemeJson)",
    "/t/[slug]/kb  — Knowledge base manager",
    "/t/[slug]/tickets  — Ticket CRUD",
    "/t/[slug]/chat  — RAG chat (client component)",
    "/api/chat  — BFF proxy; bearer token never exposed to JS",
]
txb = s.shapes.add_textbox(Inches(0.45), Inches(1.82), Inches(3.8), Inches(5.2))
tf  = txb.text_frame; tf.word_wrap = True
first = True
for text in pages:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(5)
    r = p.add_run(); r.text = "• " + text
    r.font.size = Pt(9.8); r.font.color.rgb = DARK_TEXT; r.font.name = "Segoe UI"

content_box(s, Inches(4.6), Inches(1.3), Inches(4.4), Inches(2.9))
add_text_box(s, "Key Patterns", Inches(4.75), Inches(1.38), Inches(4.1), Inches(0.4),
             font_size=Pt(13), bold=True, color=NAVY)
patterns = [
    "middleware.ts — Edge tenant resolution",
    "  Reads Host header → /t/{slug} → sets x-resolved-tenant-slug",
    "lib/tenant.ts — resolveTenant() with React cache()",
    "  One render pass = one tenant-api call (deduplicated)",
    "lib/api.ts — apiFetch wrapper",
    "  Reads httpOnly cookie; never exposes token to browser",
]
txb = s.shapes.add_textbox(Inches(4.75), Inches(1.82), Inches(4.1), Inches(2.3))
tf  = txb.text_frame; tf.word_wrap = True
first = True
for t in patterns:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(3)
    r = p.add_run(); r.text = t
    r.font.size = Pt(9.5)
    r.font.color.rgb = DARK_TEXT if not t.startswith("  ") else RGBColor(0x55, 0x55, 0x55)
    r.font.italic = t.startswith("  ")
    r.font.name = "Segoe UI"

content_box(s, Inches(4.6), Inches(4.35), Inches(4.4), Inches(2.85), fill=RGBColor(0xff, 0xf0, 0xe0))
add_text_box(s, "Security Design", Inches(4.75), Inches(4.43), Inches(4.1), Inches(0.4),
             font_size=Pt(12), bold=True, color=ORANGE)
sec = [
    "httpOnly JWT cookie — XSS cannot steal the token",
    "BFF /api/chat — bearer token never reaches the browser",
    "Server Actions for mutations — CSRF protected",
    "No x-tenant-id set by client — server is the authority",
]
txb = s.shapes.add_textbox(Inches(4.75), Inches(4.88), Inches(4.1), Inches(2.2))
tf  = txb.text_frame; tf.word_wrap = True
first = True
for t in sec:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(5)
    r = p.add_run(); r.text = "✓  " + t
    r.font.size = Pt(10); r.font.color.rgb = RGBColor(0x5a, 0x30, 0x00); r.font.name = "Segoe UI"

content_box(s, Inches(9.15), Inches(1.3), Inches(3.85), Inches(5.9))
add_text_box(s, "Demo Steps", Inches(9.3), Inches(1.38), Inches(3.5), Inches(0.4),
             font_size=Pt(13), bold=True, color=NAVY)
steps = [
    ("1", "Open localhost:3000 — home page with two tenant cards"),
    ("2", "Click /t/contoso vs /t/fabrikam — different branding, same codebase"),
    ("3", "DevTools → Network — page arrives fully rendered (no loading spinner)"),
    ("4", "Login → DevTools → Application → Cookies → tenantops_token is httpOnly"),
    ("5", "Open /t/contoso/kb, /tickets, /chat — show tenant-scoped data"),
    ("6", "Point to /api/chat route.ts — BFF proxies to AI; browser never sees bearer"),
]
y_off = Inches(1.82)
for num, text in steps:
    add_rect(s, Inches(9.3), y_off, Inches(0.3), Inches(0.3), fill=AZURE)
    add_text_box(s, num, Inches(9.3), y_off + Inches(0.02), Inches(0.3), Inches(0.28),
                 font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, text, Inches(9.65), y_off, Inches(3.25), Inches(0.55),
                 font_size=Pt(9.5), color=DARK_TEXT)
    y_off += Inches(0.88)

add_notes(s, """SLIDE 4 — CLIENT LAYER (3-4 minutes)

DEMO STEPS:

Step 1: Open localhost:3000 in Chrome. Show the two tenant cards (Contoso HVAC, Fabrikam Services).
SAY: "These are two completely separate tenants sharing the same codebase and infrastructure."

Step 2: Click into /t/contoso. Point out the Contoso branding (colors, logo, company name). Go back, click /t/fabrikam. Different branding, same code.
SAY: "The branding comes from ThemeJson stored per tenant in the database. The middleware.ts file reads the hostname or slug and resolves the tenant before the page renders."

Step 3: Right-click → Inspect → Network. Reload the page.
SAY: "Notice this is a React Server Component — the HTML arrives pre-rendered with the tenant data already embedded. There's no client-side fetch waterfall."

Step 4: Application tab → Cookies → show tenantops_token.
SAY: "The JWT is stored as an httpOnly cookie. This means no JavaScript on the page can read it. XSS attacks cannot steal the token. The BFF pattern in /api/chat reads it server-side and forwards it as a bearer token to the AI service — the browser never sees the raw token."

Step 5: Navigate to /t/contoso/tickets. Show real ticket data. Navigate to /t/fabrikam/tickets. Different data.

KEY TALKING POINT: "The client layer never sets x-tenant-id directly. The tenant identity flows exclusively from the server-validated JWT cookie. Even if a user modifies their localStorage or forges a header in DevTools, the server ignores it." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Layer 2: Edge / CDN / WAF
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 2 — Edge / CDN / WAF",
           "Azure Front Door Premium  ·  WAF Prevention Mode  ·  OWASP Managed Rules  ·  Bot Manager")
layer_badge(s, "EDGE / WAF LAYER")

flow_boxes = [
    ("Internet\nUser", Inches(0.35), Inches(2.0), Inches(1.4), Inches(1.1), NAVY),
    ("Azure\nFront Door\nPremium", Inches(2.05), Inches(1.85), Inches(1.8), Inches(1.4), RGBColor(0x00, 0x62, 0xb1)),
    ("WAF\nPolicy", Inches(4.15), Inches(1.85), Inches(1.5), Inches(1.4), RED),
    ("Route /*\nWeb App\n(Container Apps)", Inches(6.0), Inches(1.65), Inches(2.0), Inches(0.9), GREEN),
    ("Route /api/*\nAPIM", Inches(6.0), Inches(2.8), Inches(2.0), Inches(0.9), PURPLE),
]
for label, x, y, w, h, col in flow_boxes:
    add_rect(s, x, y, w, h, fill=col)
    add_text_box(s, label, x + Inches(0.05), y + Inches(0.1),
                 w - Inches(0.1), h - Inches(0.1),
                 font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ax, ay, aw, ah in [
    (Inches(1.75), Inches(2.5), Inches(0.3), Inches(0.05)),
    (Inches(3.85), Inches(2.3), Inches(0.3), Inches(0.05)),
    (Inches(5.65), Inches(2.05), Inches(0.35), Inches(0.05)),
    (Inches(5.65), Inches(3.2), Inches(0.35), Inches(0.05)),
]:
    add_rect(s, ax, ay, aw, ah, fill=NAVY)

content_box(s, Inches(0.3), Inches(3.5), Inches(8.0), Inches(3.7))
add_text_box(s, "WAF Configuration  (infra/terraform/modules/frontdoor/main.tf)",
             Inches(0.45), Inches(3.58), Inches(7.7), Inches(0.42),
             font_size=Pt(12), bold=True, color=NAVY)
waf_rows = [
    ("Setting", "Value", "Why It Matters", True),
    ("SKU", "Front Door Premium", "Required for WAF managed rules", False),
    ("WAF Mode", "Prevention", "Blocks (not just logs) malicious requests", False),
    ("Managed Rule Set", "DefaultRuleSet v2.1", "OWASP Top 10 coverage", False),
    ("Bot Manager", "BotManager v1.1", "Blocks known bad bots and scrapers", False),
    ("Custom Rate Limit", "600 req/min per IP", "Edge-level DDoS mitigation before APIM", False),
    ("HTTPS Redirect", "Enabled on /* route", "All HTTP traffic upgraded at the edge", False),
    ("Origin", "Internal Container Apps only", "Services never directly internet-accessible", False),
]
col_ws = [Inches(1.8), Inches(2.4), Inches(3.5)]
col_xs = [Inches(0.45), Inches(2.3), Inches(4.75)]
ry = Inches(4.05)
for row in waf_rows:
    rh = Inches(0.3) if row[3] else Inches(0.34)
    for ci, (cx, cw, cell) in enumerate(zip(col_xs, col_ws, row[:3])):
        bg = NAVY if row[3] else (LIGHT_BLUE if waf_rows.index(row) % 2 == 1 else WHITE)
        add_rect(s, cx - Inches(0.05), ry, cw + Inches(0.05), rh,
                 fill=bg, line=MID_GREY, line_w=Pt(0.5))
        add_text_box(s, cell, cx, ry + Inches(0.05), cw, rh - Inches(0.06),
                     font_size=Pt(9.5), bold=row[3],
                     color=WHITE if row[3] else DARK_TEXT)
    ry += rh

content_box(s, Inches(8.5), Inches(1.3), Inches(4.5), Inches(5.9))
add_text_box(s, "Local Dev Equivalent", Inches(8.65), Inches(1.38), Inches(4.2), Inches(0.4),
             font_size=Pt(13), bold=True, color=NAVY)
add_text_box(s,
    "Front Door does not run locally.\nServices are called directly via docker compose.\n\nSimulate rate-limit concept:",
    Inches(8.65), Inches(1.82), Inches(4.2), Inches(1.0), font_size=Pt(10), color=DARK_TEXT)
code_block(s, [
    "# Rapid requests (WAF would",
    "# block after 600/min in prod)",
    "for i in {1..10}; do",
    "  curl -s -o /dev/null \\",
    "    -w \"%{http_code}\\n\" \\",
    "    http://localhost:3000/t/contoso",
    "done",
    "# Production: 429 after threshold",
], Inches(8.5), Inches(2.95), Inches(4.5), Inches(2.3))
add_text_box(s, "Key Talking Point",
             Inches(8.65), Inches(5.35), Inches(4.2), Inches(0.38),
             font_size=Pt(11), bold=True, color=ORANGE)
add_text_box(s,
    "Front Door terminates TLS and routes /api/* to APIM. Origin services "
    "are internal-only — never directly reachable from the internet.",
    Inches(8.65), Inches(5.72), Inches(4.2), Inches(1.35),
    font_size=Pt(10), color=DARK_TEXT, italic=True)

add_notes(s, """SLIDE 5 — EDGE / WAF LAYER (2-3 minutes)

SAY: "This layer doesn't run in docker compose locally — it's a cloud-only Azure service. But let's walk through what it does in production."

OPEN: infra/terraform/modules/frontdoor/main.tf
Point to the azurerm_cdn_frontdoor_firewall_policy resource block.
SAY: "WAF mode is Prevention, not Detection. That means it doesn't just log suspicious requests — it blocks them before they ever reach APIM or any backend service."

WALK THROUGH THE TABLE:
- OWASP 2.1 managed rule set covers SQL injection, XSS, path traversal, CSRF, and 7 other categories automatically.
- Bot Manager 1.1 uses Microsoft's threat intelligence to block known bad actors.
- The 600 req/min rate limit is per IP, applied at the edge — so DDoS traffic is cut before APIM even processes a single JWT.

ROUTING: "Front Door has exactly two routes:
- /* sends web traffic to the Next.js Container App.
- /api/* sends API traffic to APIM."

ORIGIN RESTRICTION: "The Container Apps are configured to only accept traffic from Front Door's managed service tag. You cannot reach localhost:5001 from the public internet in the Azure deployment — only from Front Door."

LOCAL DEMO: Run the curl loop. Show all 200s locally. SAY: 'In production, after the 601st request in a minute from this IP, Front Door returns 429 — the backend never sees the excess traffic.' """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Layer 3: API Gateway
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 3 — API Gateway (APIM)",
           "JWT validation  ·  Tenant ID injection  ·  Dual-scheme auth  ·  Per-tenant rate limiting")
layer_badge(s, "API GATEWAY LAYER")

add_text_box(s, "Authenticated API Policy Flow  (infra/apim/policies/api-authenticated.xml)",
             Inches(0.3), Inches(1.3), Inches(8.0), Inches(0.38),
             font_size=Pt(12), bold=True, color=NAVY)
policy_steps = [
    ("Strip\nx-tenant-id", RED),
    ("Validate\nJWT\n(Entra or Local)", RGBColor(0x00, 0x62, 0xb1)),
    ("Check\ntid_app\nclaim exists", RGBColor(0x00, 0x62, 0xb1)),
    ("Inject\nx-tenant-id\nfrom JWT", GREEN),
    ("Rate limit\n60/min\nper tenant", PURPLE),
    ("Forward\nto service", NAVY),
]
bw2 = Inches(1.85); bh2 = Inches(1.05); by2 = Inches(1.78); gap2 = Inches(0.12); sx2 = Inches(0.3)
for i, (label, col) in enumerate(policy_steps):
    bx2 = sx2 + i * (bw2 + gap2)
    add_rect(s, bx2, by2, bw2, bh2, fill=col)
    add_text_box(s, label, bx2 + Inches(0.05), by2 + Inches(0.1),
                 bw2 - Inches(0.1), bh2 - Inches(0.1),
                 font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(policy_steps) - 1:
        add_rect(s, bx2 + bw2, by2 + bh2 / 2 - Inches(0.02), gap2, Inches(0.05), fill=NAVY)

add_text_box(s, "Policy highlights — why this order matters:",
             Inches(0.3), Inches(3.1), Inches(6.5), Inches(0.38),
             font_size=Pt(11), bold=True, color=NAVY)
code_block(s, [
    "<!-- 1. Client CANNOT set x-tenant-id -->",
    "<set-header name=\"x-tenant-id\" exists-action=\"delete\"/>",
    "",
    "<!-- 2. Dual JWT: Entra OR local dev key -->",
    "<choose>",
    "  <when condition=\"@(jwt.Claims['iss']...Entra...)\">",
    "    <validate-jwt openid-config=\"{entra-oidc-url}\"/>",
    "  </when>",
    "  <otherwise>",
    "    <validate-jwt key=\"{local-sym-key}\"/>",
    "  </otherwise>",
    "</choose>",
    "",
    "<!-- 3. Inject authoritative tenant header -->",
    "<set-header name=\"x-tenant-id\">",
    "  <value>@(context.Request.Claims[\"tid_app\"].First())</value>",
    "</set-header>",
    "",
    "<!-- 4. Per-tenant rate limit: 60/min -->",
    "<rate-limit-by-key calls=\"60\" renewal-period=\"60\"",
    "  counter-key=\"@(context.Request.Claims[\"tid_app\"].First())\"/>",
], Inches(0.3), Inches(3.55), Inches(6.4), Inches(3.7))

content_box(s, Inches(6.85), Inches(1.3), Inches(6.15), Inches(2.45))
add_text_box(s, "4 APIs registered in APIM", Inches(7.0), Inches(1.38), Inches(5.8), Inches(0.4),
             font_size=Pt(12), bold=True, color=NAVY)
apis = [
    ("/tenant  →", "tenant-api  (port 5001)", "resolve is public; rest authenticated"),
    ("/core    →", "core-api    (port 5002)", "tickets + documents, all authenticated"),
    ("/ai      →", "ai-orchestrator (5003)",  "chat + indexing, all authenticated"),
    ("/identity →","identity-api (port 5004)","login / register (no APIM needed)"),
]
ay2 = Inches(1.85)
for prefix, target, note in apis:
    add_text_box(s, prefix, Inches(7.0), ay2, Inches(1.1), Inches(0.3),
                 font_size=Pt(10), bold=True, color=AZURE, font_name="Cascadia Code")
    add_text_box(s, target, Inches(8.1), ay2, Inches(2.0), Inches(0.3),
                 font_size=Pt(10), color=DARK_TEXT, font_name="Cascadia Code")
    add_text_box(s, note, Inches(7.0), ay2 + Inches(0.28), Inches(5.8), Inches(0.25),
                 font_size=Pt(9), color=RGBColor(0x44, 0x44, 0x44), italic=True)
    ay2 += Inches(0.55)

content_box(s, Inches(6.85), Inches(3.85), Inches(6.15), Inches(2.0), fill=RGBColor(0xe8, 0xf8, 0xe8))
add_text_box(s, "Demo: Public vs Protected Endpoint",
             Inches(7.0), Inches(3.93), Inches(5.8), Inches(0.4),
             font_size=Pt(12), bold=True, color=GREEN)
code_block(s, [
    "# Public — no token (api-resolve-public.xml)",
    "curl http://localhost:5001/resolve?slug=contoso",
    "# → 200 {tenantId, slug, name, themeJson}",
    "",
    "# Protected — no token",
    "curl http://localhost:5001/me/tenant",
    "# → 401",
], Inches(6.85), Inches(4.38), Inches(6.15), Inches(1.42))
content_box(s, Inches(6.85), Inches(5.95), Inches(6.15), Inches(1.25),
            fill=RGBColor(0xff, 0xf0, 0xe0))
add_text_box(s, "Key Talking Point", Inches(7.0), Inches(6.03), Inches(5.8), Inches(0.35),
             font_size=Pt(11), bold=True, color=ORANGE)
add_text_box(s,
    "Even if a developer forgets [Authorize] on a .NET endpoint, APIM's JWT "
    "validation blocks the request before it reaches the service.",
    Inches(7.0), Inches(6.4), Inches(5.8), Inches(0.75),
    font_size=Pt(10), color=DARK_TEXT, italic=True)

add_notes(s, """SLIDE 6 — API GATEWAY / APIM (3-4 minutes)

DEMO: Run both curl commands from the code block.
SAY: "The /resolve endpoint is intentionally public — the web app needs to call it before the user is authenticated to get tenant branding. Everything else requires a valid JWT."

WALK THROUGH POLICY STEPS:
Step 1 — STRIP: "We delete any incoming x-tenant-id header. The client cannot forge their tenant. Even if a Contoso user adds 'x-tenant-id: fabrikam-guid' to their request, this header is removed before any backend service sees it."

Step 2 — VALIDATE JWT: "We support two issuers. In production, Entra ID JWTs are validated against the OIDC discovery endpoint. In local dev, identity-api signs JWTs with a symmetric key. The policy uses a <choose> block to handle both cases — same security guarantee, different key sources."

Step 3 — INJECT: "After validation, we read the tid_app claim from the validated JWT and inject it as x-tenant-id. This header is now authoritative — it came from a validated token, not from the client."

Step 4 — RATE LIMIT: "Per-tenant rate limiting. If Contoso's API key is compromised, only Contoso is throttled — Fabrikam traffic is unaffected."

KEY POINT: "This is why APIM is Layer 3 but acts as a second gatekeeper. Even if someone bypasses WAF (Layer 2), APIM still validates the JWT. Even if someone bypasses APIM (very difficult), TenantResolutionMiddleware (Layer 5) still validates the JWT. Defense in depth." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Layer 4: Application Layer
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 4 — Application Layer",
           "4 × ASP.NET Core 8 microservices  ·  ServiceDefaults bootstrap  ·  Minimal API style")
layer_badge(s, "APPLICATION LAYER")

services_data = [
    ("identity-api", ":5004", AZURE, [
        "POST /auth/register — create user",
        "POST /auth/login    — issue JWT",
        "JWT: tid_app, sub, email, name",
        "PBKDF2 password hash (100k iter)",
    ]),
    ("tenant-api", ":5001", RGBColor(0x00, 0x4c, 0x8c), [
        "GET /resolve?slug= — public (anon)",
        "GET /admin/tenants — platform-admin",
        "GET /me/tenant     — auth required",
        "5-min MemoryCache on resolve",
    ]),
    ("core-api", ":5002", GREEN, [
        "GET/POST/PUT/DELETE /tickets",
        "GET/POST/DELETE     /documents",
        "No WHERE TenantId — RLS filters",
        "TenantContext injected via DI",
    ]),
    ("ai-orchestrator", ":5003", RED, [
        "POST /chat            — RAG answer",
        "POST /documents/index — batch index",
        "POST /documents/{id}/index",
        "Semantic Kernel + Azure OpenAI",
    ]),
]
bx_start = Inches(0.25); bw3 = Inches(3.15); by3 = Inches(1.3); bh3 = Inches(5.9)
for i, (name, port, col, endpoints) in enumerate(services_data):
    bx3 = bx_start + i * (bw3 + Inches(0.1))
    add_rect(s, bx3, by3, bw3, Inches(0.65), fill=col)
    add_text_box(s, name, bx3 + Inches(0.1), by3 + Inches(0.05),
                 bw3 - Inches(0.2), Inches(0.35), font_size=Pt(13), bold=True, color=WHITE)
    add_text_box(s, "localhost" + port, bx3 + Inches(0.1), by3 + Inches(0.38),
                 bw3 - Inches(0.2), Inches(0.25),
                 font_size=Pt(10), color=RGBColor(0xcc, 0xee, 0xff), font_name="Cascadia Code")
    content_box(s, bx3, by3 + Inches(0.65), bw3, bh3 - Inches(0.65), fill=WHITE, line=col)
    txb = s.shapes.add_textbox(bx3 + Inches(0.1), by3 + Inches(0.75),
                                bw3 - Inches(0.2), bh3 - Inches(0.85))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for ep in endpoints:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(6)
        r = p.add_run(); r.text = ep
        r.font.size = Pt(10)
        r.font.name = "Cascadia Code" if ep.startswith(("GET","POST","PUT","DELETE")) else "Segoe UI"
        r.font.color.rgb = DARK_TEXT

add_rect(s, Inches(0.25), Inches(6.5), Inches(12.85), Inches(0.88), fill=CODE_BG)
add_text_box(s,
    "TOKEN=$(curl -s -X POST http://localhost:5004/auth/login -H 'Content-Type: application/json' "
    "-d '{\"tenantId\":\"11111111-1111-1111-1111-111111111111\",\"email\":\"demo@contoso.test\","
    "\"password\":\"Demo1234!\"}' | python -c \"import sys,json; print(json.load(sys.stdin)['access_token'])\")\n"
    "curl -H \"Authorization: Bearer $TOKEN\" http://localhost:5002/tickets",
    Inches(0.4), Inches(6.52), Inches(12.5), Inches(0.8),
    font_size=Pt(9), color=CODE_FG, font_name="Cascadia Code")

add_notes(s, """SLIDE 7 — APPLICATION LAYER (3-4 minutes)

DEMO: Run the two-liner at the bottom of the slide.
SAY: "Let me show what a real authenticated request looks like. We get a JWT from identity-api, then use it to call core-api for tickets."

WALK THROUGH EACH SERVICE:

identity-api: "This is the only service that issues JWTs. It does PBKDF2 password hashing with 100,000 iterations — the same algorithm recommended by NIST. The JWT payload includes tid_app (our tenant ID claim), sub (user ID), email, and name. The tenant claim is baked into the token at login time and cannot be changed without re-authenticating."

tenant-api: "The /resolve endpoint is the public lookup that maps a slug like 'contoso' to a tenant ID, theme, and branding. It has a 5-minute MemoryCache — calling it repeatedly doesn't hit the database every time. /admin/tenants requires platform-admin role."

core-api: "Notice: there is NO WHERE TenantId = @tenantId in the SQL queries. The TenantContext object is injected via DI, and the SQL connection factory uses it to set SESSION_CONTEXT before the first query. SQL RLS handles filtering automatically."

ai-orchestrator: "This service orchestrates the full RAG pipeline using Semantic Kernel. /documents/index reads documents from core-api (already RLS-filtered), embeds them, and uploads to AI Search with a tenantId field stamped on each chunk."

KEY POINT: "Each service registers with ServiceDefaults — a shared setup class that configures Serilog, auth, health checks, Swagger, and the TenantResolutionMiddleware. You add one line to Program.cs and get all security infrastructure automatically." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Layer 5: Service / Module Layer
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 5 — Service / Module Layer",
           "Shared library  ·  TenantResolutionMiddleware  ·  Dual auth schemes  ·  RLS context factory")
layer_badge(s, "SERVICE LAYER")

add_text_box(s, "Request flow through TenantResolutionMiddleware",
             Inches(0.3), Inches(1.3), Inches(8.5), Inches(0.38),
             font_size=Pt(12), bold=True, color=NAVY)
mw_steps = [
    ("Skip\n/health\n/auth\n/resolve", RGBColor(0x44, 0x44, 0x44)),
    ("Check\nIsAuthenticated\n→ 401 if not", RED),
    ("Extract\ntid_app\nclaim → 403\nif missing", RED),
    ("Compare\ntid_app vs\nx-tenant-id\n→ 403 mismatch", RED),
    ("Populate\nTenantContext\n(immutable)", RGBColor(0x00, 0x62, 0xb1)),
    ("Push to\nSerilog scope\n(TenantId\n+ UserId)", GREEN),
    ("→ Next()", GREEN),
]
bw4 = Inches(1.6); bh4 = Inches(1.1); by4 = Inches(1.78); gap4 = Inches(0.08)
for i, (label, col) in enumerate(mw_steps):
    bx4 = Inches(0.3) + i * (bw4 + gap4)
    add_rect(s, bx4, by4, bw4, bh4, fill=col)
    add_text_box(s, label, bx4 + Inches(0.06), by4 + Inches(0.08),
                 bw4 - Inches(0.12), bh4 - Inches(0.08),
                 font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(mw_steps) - 1:
        add_rect(s, bx4 + bw4, by4 + bh4 / 2 - Inches(0.02), gap4, Inches(0.04), fill=NAVY)

col_data = [
    {
        "title": "Dual Auth Scheme  (Auth/AuthSetup.cs)",
        "color": AZURE,
        "items": [
            "Policy scheme 'TenantOps' auto-selects:",
            "  Entra issuer → OIDC config validation",
            "    login.microsoftonline.com/{tenant}/v2.0",
            "  Local issuer → symmetric key (identity-api)",
            "Both require tid_app claim",
            "2-minute clock skew tolerance",
        ],
        "x": Inches(0.3), "w": Inches(4.1),
    },
    {
        "title": "TenantContext  (Tenancy/TenantContext.cs)",
        "color": PURPLE,
        "items": [
            "Request-scoped, immutable after Set()",
            "Fields: TenantId, Slug, UserId,",
            "         IsPlatformAdmin, IsSet",
            "Populated once by middleware",
            "Cannot be mutated by downstream code",
            "AI orchestrator checks IsSet before search",
        ],
        "x": Inches(4.55), "w": Inches(4.2),
    },
    {
        "title": "Connection Factory  (Data/TenantDbConnectionFactory.cs)",
        "color": GREEN,
        "items": [
            "OpenAsync():",
            "  Sets SESSION_CONTEXT(TenantId = guid)",
            "  Sets IsPlatformAdmin = 0",
            "OpenPlatformAdminAsync():",
            "  Sets IsPlatformAdmin = 1 (reads Tenants)",
            "Both: read_only=1 → context cannot change",
        ],
        "x": Inches(8.9), "w": Inches(4.1),
    },
]
for col in col_data:
    content_box(s, col["x"], Inches(3.05), col["w"], Inches(4.15))
    add_text_box(s, col["title"], col["x"] + Inches(0.15), Inches(3.13),
                 col["w"] - Inches(0.3), Inches(0.42),
                 font_size=Pt(10.5), bold=True, color=col["color"])
    txb = s.shapes.add_textbox(col["x"] + Inches(0.15), Inches(3.6),
                                col["w"] - Inches(0.3), Inches(3.5))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for item in col["items"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        r = p.add_run(); r.text = item
        r.font.size = Pt(10)
        r.font.name = "Cascadia Code" if item.startswith("  ") else "Segoe UI"
        r.font.italic = item.startswith("  ")
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

content_box(s, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.72),
            fill=RGBColor(0xff, 0xee, 0xee))
add_text_box(s, "Cross-tenant attack defeated by middleware:",
             Inches(0.45), Inches(6.65), Inches(4.0), Inches(0.35),
             font_size=Pt(11), bold=True, color=RGBColor(0x99, 0x00, 0x00))
add_text_box(s,
    "Contoso user forges x-tenant-id: fabrikam-guid  →  middleware sees tid_app=contoso-guid vs header=fabrikam-guid  →  403 immediately",
    Inches(4.5), Inches(6.65), Inches(8.4), Inches(0.6),
    font_size=Pt(10), color=DARK_TEXT, italic=True)

add_notes(s, """SLIDE 8 — SERVICE / MODULE LAYER (3-4 minutes)

OPEN: services/_shared/TenantOps.Shared/Tenancy/TenantResolutionMiddleware.cs
Walk through InvokeAsync step by step.

WALK THROUGH MIDDLEWARE STEPS:
Step 1 — SKIP: "Health checks, /auth/login, /auth/register, and /resolve are anonymous — they must work before a user is authenticated. We skip them explicitly."

Step 2 — CHECK AUTH: "If the request made it past APIM without a valid JWT, middleware catches it here with a 401."

Step 3 — EXTRACT CLAIM: "We extract the tid_app claim from the validated JWT. We also accept extension_TenantId (Entra custom attribute) and tenantId as fallbacks. If the claim is missing or malformed, we return 403."

Step 4 — COMPARE: "This is the key anti-forgery check. Even though APIM already stripped and re-injected x-tenant-id, we compare the JWT claim against the header again at the service level. If a forged header somehow survived (e.g., in local dev without APIM), this check catches it."

Step 5 — POPULATE: "TenantContext is a request-scoped DI object that becomes immutable after Set() is called. Downstream code reads TenantId from it — nobody can change it mid-request."

Step 6 — SERILOG: "Every log line for the rest of this request automatically includes TenantId and UserId. In Application Insights, you can filter all logs for a specific tenant without adding logging code to individual endpoints."

ATTACK DEMO: Point to the red box at the bottom.
SAY: "Let's say a Contoso user adds the header 'x-tenant-id: 22222222-2222-2222-2222-222222222222' to their request. Their JWT still says tid_app=11111111-... The comparison fails. 403 returned. No data leakage." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Layer 6: Data
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 6 — Data Layer",
           "Azure SQL  ·  Row-Level Security  ·  SESSION_CONTEXT  ·  Low-privilege app login")
layer_badge(s, "DATA LAYER")

schema_items = [
    ("app.Tenants", "NOT RLS-protected\nTenantId · Slug · Domain\nName · ThemeJson · IsActive\nRead via platform-admin conn",
     RGBColor(0x44, 0x44, 0x88)),
    ("app.Users", "RLS: FILTER + BLOCK\nUserId · TenantId (FK)\nEmail · DisplayName\nIdentityProvider · PasswordHash",
     RGBColor(0x00, 0x62, 0xb1)),
    ("app.Tickets", "RLS: FILTER + BLOCK\nTicketId · TenantId (FK)\nTitle · Body · Status\nCreatedBy · Timestamps",
     GREEN),
    ("app.Documents", "RLS: FILTER + BLOCK\nDocumentId · TenantId (FK)\nTitle · Content\nCreatedBy · IndexedAt",
     RED),
]
tbw = Inches(2.8); tby = Inches(1.3); tbh = Inches(2.2)
for i, (name, fields, col) in enumerate(schema_items):
    tbx = Inches(0.25) + i * (tbw + Inches(0.15))
    add_rect(s, tbx, tby, tbw, Inches(0.42), fill=col)
    add_text_box(s, name, tbx + Inches(0.1), tby + Inches(0.05),
                 tbw - Inches(0.2), Inches(0.35), font_size=Pt(11), bold=True, color=WHITE)
    content_box(s, tbx, tby + Inches(0.42), tbw, tbh - Inches(0.42), fill=WHITE, line=col)
    add_text_box(s, fields, tbx + Inches(0.1), tby + Inches(0.5),
                 tbw - Inches(0.2), tbh - Inches(0.6), font_size=Pt(9.5), color=DARK_TEXT)

add_text_box(s, "RLS Predicate Function  (002_rls.sql)",
             Inches(0.25), Inches(3.65), Inches(6.5), Inches(0.38),
             font_size=Pt(12), bold=True, color=NAVY)
code_block(s, [
    "CREATE FUNCTION app.fn_tenant_predicate(@TenantId UNIQUEIDENTIFIER)",
    "RETURNS TABLE WITH SCHEMABINDING AS",
    "RETURN SELECT 1 AS result",
    "WHERE SESSION_CONTEXT(N'TenantId') = @TenantId",
    "   OR CAST(SESSION_CONTEXT(N'IsPlatformAdmin') AS BIT) = 1;",
    "",
    "-- FILTER (SELECT) + BLOCK (INSERT/UPDATE/DELETE)",
    "CREATE SECURITY POLICY app.TenantIsolationPolicy",
    "  ADD FILTER PREDICATE app.fn_tenant_predicate(TenantId) ON app.Tickets,",
    "  ADD BLOCK  PREDICATE app.fn_tenant_predicate(TenantId) ON app.Tickets AFTER INSERT,",
    "  ADD FILTER PREDICATE app.fn_tenant_predicate(TenantId) ON app.Documents,",
    "  ADD FILTER PREDICATE app.fn_tenant_predicate(TenantId) ON app.Users;",
], Inches(0.25), Inches(4.08), Inches(6.6), Inches(3.2))

content_box(s, Inches(7.0), Inches(3.65), Inches(6.0), Inches(3.6))
add_text_box(s, "Live Demo — RLS in SQL console",
             Inches(7.15), Inches(3.73), Inches(5.7), Inches(0.42),
             font_size=Pt(12), bold=True, color=NAVY)
code_block(s, [
    "-- Connect as app login (not sa!)",
    "sqlcmd -U tenantops_app -P App_Strong_Passw0rd!",
    "",
    "-- Session: Contoso",
    "EXEC sp_set_session_context",
    "  @key=N'TenantId',",
    "  @value='11111111-1111-1111-1111-111111111111';",
    "SELECT Title FROM app.Documents;",
    "-- → Contoso docs only (no WHERE needed)",
    "",
    "-- Cross-tenant INSERT attempt",
    "INSERT INTO app.Documents (TenantId, Title, Content)",
    "VALUES ('22222222-...', 'hack', 'x');",
    "-- → Error: BLOCK predicate blocked INSERT",
], Inches(7.0), Inches(4.18), Inches(6.0), Inches(3.0))

add_notes(s, """SLIDE 9 — DATA LAYER (4-5 minutes) — MOST IMPACTFUL DEMO

This is the most visually compelling demo. Live SQL console.

SETUP: Open a new terminal.
RUN: docker compose exec sql /opt/mssql-tools/bin/sqlcmd -S localhost -U tenantops_app -P "App_Strong_Passw0rd!" -d tenantops

STEP 1 — Set Contoso context and SELECT:
EXEC sp_set_session_context @key=N'TenantId', @value='11111111-1111-1111-1111-111111111111', @read_only=1;
SELECT Title FROM app.Documents;

SAY: "We get only Contoso documents. There is no WHERE clause in this query. The database engine automatically applies the RLS predicate. The application developer doesn't need to remember to filter — the database enforces it."

STEP 2 — Try cross-tenant INSERT:
INSERT INTO app.Documents (TenantId, Title, Content) VALUES ('22222222-2222-2222-2222-222222222222', 'Injected doc', 'malicious content');

SAY: "403 from the database. The BLOCK predicate prevents inserting a row for Fabrikam while the session context is set to Contoso. Even if an attacker got a SQL connection, they cannot write to another tenant's data."

STEP 3 — Switch to Fabrikam context:
EXEC sp_set_session_context @key=N'TenantId', @value='22222222-2222-2222-2222-222222222222', @read_only=1;
SELECT Title FROM app.Documents;

SAY: "Completely different data. Same connection, same query, different SESSION_CONTEXT."

EXPLAIN THE SCHEMA: app.Tenants is NOT RLS-protected because the platform admin needs to read all tenants. It's accessed via a separate connection (OpenPlatformAdminAsync) with IsPlatformAdmin=1 in SESSION_CONTEXT.

KEY POINT: "The app login (tenantops_app) has minimal permissions — INSERT/UPDATE/DELETE/SELECT on the app schema tables only. Even if an attacker gained database credentials, they still can't bypass RLS." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Layer 7: Infrastructure
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 7 — Infrastructure Layer",
           "Docker Compose (local)  ·  Terraform IaC (Azure)  ·  Managed identities  ·  No secrets in code")
layer_badge(s, "INFRA LAYER")

content_box(s, Inches(0.25), Inches(1.3), Inches(5.6), Inches(5.9))
add_text_box(s, "Local Stack  (docker-compose.yml)",
             Inches(0.4), Inches(1.38), Inches(5.3), Inches(0.42),
             font_size=Pt(12), bold=True, color=NAVY)
svc_rows = [
    ("tenantops-sql",          "Azure SQL Edge",  "Port 1433   sqlcmd healthcheck"),
    ("tenantops-azurite",      "Azure Storage",   "Ports 10000–10002 (blob/queue/table)"),
    ("tenantops-identity-api", ".NET 8 service",  "Port 5004   depends: sql"),
    ("tenantops-tenant-api",   ".NET 8 service",  "Port 5001   depends: sql"),
    ("tenantops-core-api",     ".NET 8 service",  "Port 5002   depends: sql"),
    ("tenantops-ai",           ".NET 8 service",  "Port 5003   depends: sql"),
    ("tenantops-web",          "Next.js 14",      "Port 3000   depends: all services"),
]
ry2 = Inches(1.88)
for name, img, note in svc_rows:
    add_rect(s, Inches(0.4), ry2, Inches(2.2), Inches(0.38), fill=NAVY)
    add_text_box(s, name, Inches(0.45), ry2 + Inches(0.03), Inches(2.1), Inches(0.32),
                 font_size=Pt(8.5), color=WHITE, font_name="Cascadia Code")
    add_text_box(s, img, Inches(2.65), ry2 + Inches(0.03), Inches(1.3), Inches(0.32),
                 font_size=Pt(8.5), bold=True, color=AZURE)
    add_text_box(s, note, Inches(2.65), ry2 + Inches(0.28), Inches(3.1), Inches(0.25),
                 font_size=Pt(8), color=RGBColor(0x55, 0x55, 0x55), italic=True)
    ry2 += Inches(0.74)
add_text_box(s, "make up  →  auto-migrate + seed",
             Inches(0.4), Inches(6.82), Inches(5.3), Inches(0.32),
             font_size=Pt(10), color=AZURE, italic=True, font_name="Cascadia Code")

content_box(s, Inches(6.0), Inches(1.3), Inches(7.05), Inches(3.5))
add_text_box(s, "Terraform Module Composition  (infra/terraform/)",
             Inches(6.15), Inches(1.38), Inches(6.7), Inches(0.42),
             font_size=Pt(12), bold=True, color=NAVY)
tf_items = [
    ("main.tf", "Resource group · Log Analytics · App Insights · ACR · Managed Identity"),
    ("modules/sql", "Azure SQL Server + Database · Entra admin · audit logs → LA"),
    ("modules/ai", "Azure OpenAI (S0) + AI Search · RBAC roles for app identity"),
    ("modules/aca", "Container Apps env · 4 apps · Dapr sidecar · internal ingress"),
    ("modules/apim", "API Management · 4 APIs · global + auth + public policies"),
    ("modules/frontdoor", "Front Door Premium · WAF · 2 routes · health probes"),
]
txb = s.shapes.add_textbox(Inches(6.15), Inches(1.88), Inches(6.7), Inches(2.8))
tf2 = txb.text_frame; tf2.word_wrap = True
first = True
for mod, desc in tf_items:
    p = tf2.paragraphs[0] if first else tf2.add_paragraph()
    first = False
    p.space_before = Pt(5)
    r1 = p.add_run(); r1.text = f"{mod}   "
    r1.font.size = Pt(10); r1.font.bold = True; r1.font.name = "Cascadia Code"
    r1.font.color.rgb = AZURE
    r2 = p.add_run(); r2.text = desc
    r2.font.size = Pt(10); r2.font.name = "Segoe UI"; r2.font.color.rgb = DARK_TEXT

content_box(s, Inches(6.0), Inches(4.95), Inches(7.05), Inches(2.25))
add_text_box(s, "Local ↔ Azure Parity",
             Inches(6.15), Inches(5.03), Inches(6.7), Inches(0.42),
             font_size=Pt(12), bold=True, color=NAVY)
parity = [
    ("Concern", "Local Dev", "Azure Production", True),
    ("SQL",      "Azure SQL Edge container", "Azure SQL Database (managed)", False),
    ("Auth",     "Symmetric JWT key",       "Entra ID OIDC", False),
    ("Gateway",  "Direct service calls",    "APIM + Front Door", False),
    ("Identity", "Password in .env",        "Managed Identity (no secrets)", False),
    ("TLS",      "HTTP",                    "HTTPS (Front Door terminates)", False),
]
py2 = Inches(5.5); ph = Inches(0.3)
col3_xs = [Inches(6.1), Inches(7.75), Inches(9.85)]
col3_ws = [Inches(1.55), Inches(2.05), Inches(3.1)]
for row in parity:
    for ci, (cx, cw, cell) in enumerate(zip(col3_xs, col3_ws, row[:3])):
        bg = NAVY if row[3] else (LIGHT_BLUE if parity.index(row) % 2 == 0 else WHITE)
        add_rect(s, cx, py2, cw, ph, fill=bg, line=MID_GREY, line_w=Pt(0.5))
        add_text_box(s, cell, cx + Inches(0.06), py2 + Inches(0.04), cw - Inches(0.1), ph,
                     font_size=Pt(9), bold=row[3], color=WHITE if row[3] else DARK_TEXT)
    py2 += ph

add_notes(s, """SLIDE 10 — INFRASTRUCTURE LAYER (2-3 minutes)

DEMO: Run 'docker compose ps' in the terminal.
SAY: "We have 14 containers running. The key design decision was making local dev as close to Azure production as possible."

WALK THROUGH LOCAL STACK:
- tenantops-sql runs Azure SQL Edge — same SQL Server engine, same T-SQL, same RLS. No schema changes needed when deploying to Azure.
- tenantops-azurite emulates Azure Blob, Queue, and Table storage locally.
- All .NET services depend on the SQL container's healthcheck before starting.

WALK THROUGH TERRAFORM MODULES: Open infra/terraform/main.tf
SAY: "The Terraform modules compose the full Azure deployment. The key property: no secrets in the Terraform code."

POINT TO MANAGED IDENTITY: "In Azure, the Container Apps environment is assigned a managed identity. The sql module grants that identity the 'db_datareader' and 'db_datawriter' roles on the SQL database. The AI module grants it the 'Cognitive Services OpenAI User' role on Azure OpenAI and the 'Search Index Data Contributor' role on AI Search. No connection strings, no API keys in environment variables."

POINT OUT THE PARITY TABLE: "The code works identically in both environments. The only difference is which JWT issuer the auth scheme accepts. This is toggled by an environment variable, not a code change."

KEY POINT: "make up runs docker compose up, waits 15 seconds, then calls make migrate and make seed automatically. One command brings up the entire stack with test data." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Layer 8: AI Layer
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Layer 8 — AI Layer",
           "Semantic Kernel  ·  Azure OpenAI  ·  Azure AI Search  ·  RAG with tenant-isolated vector search")
layer_badge(s, "AI LAYER")

add_text_box(s, "RAG Pipeline  (services/ai-orchestrator/Rag/RagOrchestrator.cs)",
             Inches(0.25), Inches(1.3), Inches(12.8), Inches(0.38),
             font_size=Pt(12), bold=True, color=NAVY)
rag_steps = [
    ("1. Embed\nQuery", "text-embedding\n-3-small\n→ 1536-dim\nvector", PURPLE),
    ("2. Vector\nSearch", "AI Search\nhybrid BM25\n+ KNN top-5\n+ tenant filter", RGBColor(0x00, 0x62, 0xb1)),
    ("3. Build\nPrompt", "System msg\n+ ranked\ncontext blocks\n+ citations", RGBColor(0x00, 0x4c, 0x8c)),
    ("4. Chat\nCompletion", "gpt-4o-mini\ngrounded to\ncontext only\n(no hallucination)", NAVY),
    ("5. Extract\nCitations", "Structured\n{answer,\ncitations[]}\nwith snippets", GREEN),
]
bw5 = Inches(2.2); bh5 = Inches(1.65); by5 = Inches(1.78); gap5 = Inches(0.13); sx5 = Inches(0.25)
for i, (title, detail, col) in enumerate(rag_steps):
    bx5 = sx5 + i * (bw5 + gap5)
    add_rect(s, bx5, by5, bw5, Inches(0.48), fill=col)
    add_text_box(s, title, bx5 + Inches(0.08), by5 + Inches(0.07),
                 bw5 - Inches(0.16), Inches(0.38),
                 font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    content_box(s, bx5, by5 + Inches(0.48), bw5, bh5 - Inches(0.48), fill=WHITE, line=col)
    add_text_box(s, detail, bx5 + Inches(0.1), by5 + Inches(0.55),
                 bw5 - Inches(0.2), bh5 - Inches(0.6),
                 font_size=Pt(10), color=DARK_TEXT, align=PP_ALIGN.CENTER)
    if i < len(rag_steps) - 1:
        add_rect(s, bx5 + bw5 + Inches(0.01), by5 + bh5 / 2,
                 gap5 - Inches(0.01), Inches(0.05), fill=NAVY)

content_box(s, Inches(0.25), Inches(3.6), Inches(5.85), Inches(3.65),
            fill=RGBColor(0xff, 0xf0, 0xe0))
add_text_box(s, "Tenant Isolation — 3 Independent Enforcement Points",
             Inches(0.4), Inches(3.68), Inches(5.55), Inches(0.42),
             font_size=Pt(12), bold=True, color=ORANGE)
iso_points = [
    ("SQL (RLS)", "OpenAsync() sets SESSION_CONTEXT(TenantId) — only the tenant's own docs are readable by the indexer"),
    ("AI Search field", "Every chunk is uploaded with tenantId field stamped; filter is set at index time"),
    ("Search query filter", "OData: tenantId eq '{TenantId}' on every search; runtime guard throws if IsSet=false"),
]
iy = Inches(4.18)
for i, (label, desc) in enumerate(iso_points):
    add_rect(s, Inches(0.4), iy, Inches(0.3), Inches(0.3),
             fill=[AZURE, GREEN, RED][i])
    add_text_box(s, str(i+1), Inches(0.4), iy + Inches(0.02), Inches(0.3), Inches(0.28),
                 font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, label, Inches(0.78), iy, Inches(5.25), Inches(0.28),
                 font_size=Pt(10), bold=True, color=DARK_TEXT)
    add_text_box(s, desc, Inches(0.78), iy + Inches(0.28), Inches(5.25), Inches(0.48),
                 font_size=Pt(9.5), color=DARK_TEXT, italic=True)
    iy += Inches(0.95)

content_box(s, Inches(6.25), Inches(3.6), Inches(6.85), Inches(3.65))
add_text_box(s, "Demo: Contoso vs Fabrikam — same question, isolated answers",
             Inches(6.4), Inches(3.68), Inches(6.6), Inches(0.42),
             font_size=Pt(11), bold=True, color=NAVY)
code_block(s, [
    "# 1. Index Contoso docs",
    "curl -X POST http://localhost:5003/documents/index \\",
    "  -H \"Authorization: Bearer $CTOKEN\"",
    "",
    "# 2. Ask a question as Contoso",
    "curl -X POST http://localhost:5003/chat \\",
    "  -H \"Authorization: Bearer $CTOKEN\" \\",
    "  -d '{\"message\":\"What are service hours?\"}'",
    "# → Contoso HVAC: Mon-Sat 8am-6pm...",
    "",
    "# 3. Same question as Fabrikam",
    "curl -X POST http://localhost:5003/chat \\",
    "  -H \"Authorization: Bearer $FTOKEN\" \\",
    "  -d '{\"message\":\"What are service hours?\"}'",
    "# → Fabrikam: 24/7 with weekend surcharges...",
], Inches(6.25), Inches(4.15), Inches(6.85), Inches(3.0))

add_notes(s, """SLIDE 11 — AI LAYER (4-5 minutes) — SHOWSTOPPER DEMO

This is the most impressive demo moment. Two tokens, one question, completely different answers.

SETUP (do this before the demo or in advance):
Get Contoso token: CTOKEN=$(curl -s -X POST http://localhost:5004/auth/login -H 'Content-Type: application/json' -d '{"tenantId":"11111111-1111-1111-1111-111111111111","email":"demo@contoso.test","password":"Demo1234!"}' | python -c "import sys,json; print(json.load(sys.stdin)['access_token'])")

Get Fabrikam token: FTOKEN=$(curl -s -X POST http://localhost:5004/auth/login -H 'Content-Type: application/json' -d '{"tenantId":"22222222-2222-2222-2222-222222222222","email":"demo@fabrikam.test","password":"Demo1234!"}' | python -c "import sys,json; print(json.load(sys.stdin)['access_token'])")

STEP 1 — Index Contoso docs:
curl -X POST http://localhost:5003/documents/index -H "Authorization: Bearer $CTOKEN"
SAY: "We're embedding Contoso's knowledge base documents and uploading them to Azure AI Search. Each chunk gets a tenantId field stamped."

STEP 2 — Ask as Contoso:
curl -X POST http://localhost:5003/chat -H "Authorization: Bearer $CTOKEN" -d '{"message":"What are the service hours?"}'
Show the answer: Contoso HVAC Mon-Sat 8am-6pm specific details.

STEP 3 — Same question as Fabrikam:
curl -X POST http://localhost:5003/chat -H "Authorization: Bearer $FTOKEN" -d '{"message":"What are the service hours?"}'
Show: Fabrikam 24/7 service different details.

SAY: "The AI model is the same. The vector space is the same. The search index is the same. But the OData filter 'tenantId eq fabrikam-guid' means the model only has access to Fabrikam's chunks when answering Fabrikam's user."

EXPLAIN THE 3 ISOLATION POINTS:
1. SQL RLS: when the indexer reads documents via core-api, only Contoso's docs come back (SESSION_CONTEXT).
2. AI Search field: every chunk uploaded has tenantId stamped — this cannot be changed after upload.
3. Query filter: every search request has 'tenantId eq ...' in the OData filter. If IsSet=false (TenantContext not populated), the code throws before touching the search index — a fail-safe guard." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — End-to-End Story
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
add_rect(s, 0, Inches(1.1), SLIDE_W, Inches(0.06), fill=ORANGE)
add_text_box(s, "End-to-End Request Story",
             Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.85),
             font_size=Pt(30), bold=True, color=WHITE)
add_text_box(s, "What happens when a user visits /t/contoso/chat and submits a question",
             Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.38),
             font_size=Pt(13), color=ORANGE, italic=True)

e2e_steps = [
    ("Browser", "User visits localhost:3000/t/contoso →\nmiddleware.ts reads /t/contoso → sets x-resolved-tenant-slug cookie"),
    ("Client RSC", "getCurrentTenant() → resolveTenant(slug) → fetches tenant-api/resolve?slug=contoso →\nContoso branding rendered server-side (primary #0b5cad, logo)"),
    ("Login", "Login form → server action → identity-api/auth/login →\nJWT issued with tid_app=11111111-... → stored as httpOnly cookie"),
    ("Chat Submit", "User types question → browser POSTs to /api/chat (BFF) →\nserver reads cookie, attaches Authorization header → proxies to ai-orchestrator"),
    ("Middleware", "TenantResolutionMiddleware: validates JWT → extracts tid_app →\ncompares to x-tenant-id header → populates TenantContext(TenantId, UserId)"),
    ("RAG", "RagOrchestrator: embeds query → SearchRetriever applies filter tenantId eq 'contoso-guid' →\nretrieves top-5 chunks → builds grounded prompt → gpt-4o-mini → citations"),
    ("Data", "SQL connection: OpenAsync() sets SESSION_CONTEXT(TenantId=contoso) →\nRLS fn_tenant_predicate ensures only Contoso rows visible to any query"),
    ("Production +", "Front Door WAF filters malicious traffic → APIM validates JWT + injects x-tenant-id →\nContainer Apps (internal only) → same middleware + RLS + AI filter stack"),
]
cols_e2e = 2; col_w_e2e = Inches(6.4)
for i, (stage, desc) in enumerate(e2e_steps):
    col = i % cols_e2e; row = i // cols_e2e
    ex = Inches(0.25) + col * (col_w_e2e + Inches(0.2))
    ey = Inches(1.35) + row * Inches(1.45)
    add_rect(s, ex, ey, Inches(1.5), Inches(0.42), fill=ORANGE)
    add_text_box(s, stage, ex + Inches(0.08), ey + Inches(0.05),
                 Inches(1.35), Inches(0.35), font_size=Pt(10), bold=True, color=WHITE)
    add_rect(s, ex + Inches(1.5), ey, col_w_e2e - Inches(1.5), Inches(0.42),
             fill=RGBColor(0x1a, 0x2e, 0x45))
    add_text_box(s, desc, ex + Inches(1.52), ey + Inches(0.05),
                 col_w_e2e - Inches(1.6), Inches(1.35),
                 font_size=Pt(9), color=MID_GREY)

add_rect(s, 0, Inches(7.1), SLIDE_W, Inches(0.06), fill=ORANGE)
add_text_box(s,
    "Every layer independently enforces tenant isolation — no single layer failure can expose another tenant's data.",
    Inches(0.4), Inches(7.17), Inches(12.5), Inches(0.3),
    font_size=Pt(10), color=ORANGE, italic=True, align=PP_ALIGN.CENTER)

add_notes(s, """SLIDE 12 — END-TO-END STORY (3-4 minutes)

SAY: "Let's trace a single request through all 8 layers. A Contoso user opens the chat page and asks 'What are the service hours?'"

WALK THROUGH EACH STEP:

Browser: "The URL /t/contoso is intercepted by Next.js Edge Middleware before any React code runs. It reads the slug from the path and sets a cookie so the server components know which tenant to render."

Client RSC: "resolveTenant() calls tenant-api's public /resolve endpoint. This is the one unauthenticated call — it just maps 'contoso' to a tenant ID and branding. React cache() ensures this happens once per render cycle."

Login: "The JWT is issued by identity-api with the tenant ID baked in as tid_app. This claim cannot be changed — you'd need to re-authenticate with that tenant's credentials."

Chat Submit: "The browser never touches the bearer token. The /api/chat BFF reads the httpOnly cookie server-side, attaches Authorization header, and proxies to ai-orchestrator. From the browser's perspective, it's a same-origin POST."

Middleware: "TenantResolutionMiddleware runs at the start of every request in every service. It does the JWT claim extraction and header comparison. TenantContext is now populated."

RAG: "RagOrchestrator.AnswerAsync() first checks TenantContext.IsSet — if false, it throws immediately. The OData filter is hardcoded from TenantContext.TenantId — the caller cannot override it."

Data: "When the Connection Factory opens a SQL connection, the first thing it does is SET SESSION_CONTEXT. All subsequent queries in that connection scope are automatically filtered."

Production +: "In Azure, add two more layers before middleware: WAF blocks malformed requests at the edge, and APIM validates JWTs and injects the authoritative x-tenant-id."

CLOSING: "Count the layers: WAF → APIM → Middleware → SQL RLS → AI Search filter. Five independent checks. Any one of them alone would prevent cross-tenant access. Together, they form a defense-in-depth architecture that is auditable, testable, and observable." """)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Defense-in-Depth Summary
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GREY)
header_bar(s, "Defense-in-Depth Summary",
           "5 independent tenant isolation checkpoints — any single layer is sufficient to prevent cross-tenant access")

isolation_layers = [
    ("Edge WAF\n(Front Door)", "600 req/min per IP\nOWASP Top 10 rules\nBot Manager v1.1\nTLS termination", RGBColor(0x00, 0x62, 0xb1)),
    ("API Gateway\n(APIM)", "JWT signature + expiry\nStrips client x-tenant-id\nInjects from JWT claim\n60 req/min per tenant", PURPLE),
    ("Middleware\n(.NET Shared)", "Claim presence check\nHeader/claim agreement\nTenantContext population\nSerilog scope injection", RED),
    ("Database\n(SQL RLS)", "SESSION_CONTEXT filter\nBLOCK on INSERT/UPDATE\nLow-privilege app login\nNo WHERE clause needed", GREEN),
    ("AI Search\n(Vector)", "tenantId field on chunks\nOData filter mandatory\nRuntime IsSet guard\nNo cross-tenant vectors", ORANGE),
]
bw6 = Inches(2.3); bh6 = Inches(4.2); by6 = Inches(1.4)
for i, (title, desc, col) in enumerate(isolation_layers):
    bx6 = Inches(0.3) + i * (bw6 + Inches(0.15))
    add_rect(s, bx6, by6, bw6, Inches(0.65), fill=col)
    add_text_box(s, title, bx6 + Inches(0.1), by6 + Inches(0.1),
                 bw6 - Inches(0.2), Inches(0.52),
                 font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    content_box(s, bx6, by6 + Inches(0.65), bw6, bh6 - Inches(0.65), fill=WHITE, line=col)
    txb = s.shapes.add_textbox(bx6 + Inches(0.15), by6 + Inches(0.8),
                                bw6 - Inches(0.3), bh6 - Inches(0.9))
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for line in desc.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(7)
        r = p.add_run(); r.text = "✓  " + line
        r.font.size = Pt(10.5); r.font.name = "Segoe UI"; r.font.color.rgb = DARK_TEXT

content_box(s, Inches(0.3), Inches(5.78), Inches(12.7), Inches(1.5),
            fill=RGBColor(0xe8, 0xf8, 0xe8))
add_text_box(s, "Observability across all layers:",
             Inches(0.5), Inches(5.86), Inches(12.3), Inches(0.38),
             font_size=Pt(12), bold=True, color=GREEN)
obs_items = [
    "Log Analytics Workspace — 30-day retention, all service logs with TenantId + UserId in every line",
    "Application Insights — distributed traces, dependency tracking (AI Search, OpenAI latency)",
    "Correlation IDs — minted at Front Door, echoed through APIM, propagated via Serilog context through every log",
    "Structured logging — filter by TenantId in Application Insights to see all activity for a single tenant",
]
obs_txb = s.shapes.add_textbox(Inches(0.5), Inches(6.28), Inches(12.3), Inches(0.9))
obs_tf = obs_txb.text_frame; obs_tf.word_wrap = True
first = True
for item in obs_items:
    p = obs_tf.paragraphs[0] if first else obs_tf.add_paragraph()
    first = False
    p.space_before = Pt(2)
    r = p.add_run(); r.text = "•  " + item
    r.font.size = Pt(10); r.font.name = "Segoe UI"; r.font.color.rgb = DARK_TEXT

add_notes(s, """SLIDE 13 — DEFENSE-IN-DEPTH SUMMARY (2-3 minutes, closing)

SAY: "Let me close with the security architecture in one view. Five independent isolation checkpoints."

WALK THROUGH EACH COLUMN:

Edge WAF: "First line. Doesn't know about tenants at all — just blocks malformed, malicious, or excessive traffic. By the time a request gets past WAF, it's at least not a known-bad request."

API Gateway: "Second line. Validates JWT cryptographically. Cannot be bypassed by header manipulation. Even if WAF goes down, APIM still rejects unsigned tokens."

Middleware: "Third line. Even if someone gets a valid JWT for the wrong tenant (nearly impossible — they'd need the tenant's credentials), middleware catches the claim/header mismatch."

Database: "Fourth line. Even if an attacker gained a database connection with the app login credentials, RLS prevents them from seeing or writing to another tenant's rows."

AI Search: "Fifth line. Even if all four previous layers failed, the OData filter on every search query means the AI model cannot retrieve another tenant's documents."

SAY: "The word 'independent' means each layer would catch a cross-tenant attack even if all others were bypassed. This is what 'defense in depth' means in practice — not one strong wall, but five independent checkpoints."

OBSERVABILITY CLOSE: "And every layer is observable. Every log line carries TenantId and UserId. Every HTTP request has a correlation ID from Front Door to Application Insights. In an incident, you can reconstruct the complete path of any request across all layers."

QUESTIONS? Key questions to anticipate:
- 'What about Kubernetes/AKS?' — Same pattern applies; APIM and RLS work the same way.
- 'What about shared database costs?' — This is the shared-database, row-level-isolation model. For highest isolation requirements, tenant-per-database is also supported by the same Terraform pattern.
- 'How do you add a new tenant?' — Insert a row in app.Tenants, run seed SQL. The RLS policy automatically covers the new TenantId." """)

# ── Save ─────────────────────────────────────────────────────────────────
out_path = r"C:\Users\v-abdulmuizb\OneDrive - Microsoft\Desktop\TenantOps_Architecture_Demo_v2.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
